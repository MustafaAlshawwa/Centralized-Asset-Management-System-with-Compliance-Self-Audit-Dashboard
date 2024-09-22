import hashlib
import os
import re
import requests
import json
import socket
from docx import Document
import openpyxl
from pptx import Presentation
import PyPDF2
from datetime import datetime, timedelta

# Regular expressions for sensitive information
credit_card_regex = r'(?:(?:4\d{3})|(?:5[1-5]\d{2})|(?:6(?:011|5\d\d))|(?:3[47]\d{2}))[-\s]?\d{4}[-\s]?\d{4}[-\s]?\d{4}(?:[-\s]?\d{3})?'
health_record_regex = r'\b(?:Patient ID|Medical Record):?\s*\d{4,}\b'
financial_info_regex = r'\b(?:Account Number:?\s*\d{6,12}|Routing Number:?\s*\d{9})\b'
intellectual_property_regex = r'(?:Patent No\.|Copyright):\s*\w+'
legal_documents_regex = r'(?:Contract|Agreement)\sNo\.:\s*\d+'
gov_id_regex = r'\b(?:SSN|Social Security No\.):?\s*\d{3}-\d{2}-\d{4}\b'
educational_records_regex = r'\b(Student ID|Grade):\s*\w+'
corporate_info_regex = r'(?:Business Plan|Strategic Objective):\s*\w+'
employee_data_regex = r'\bEmployee ID:\s*\d{5,}\b'
customer_data_regex = r'\b(Customer ID|Purchase History):\s*\w+'
communication_records_regex = r'(?:Email:?\s*\b[A-Za-z0.9._%+-]+@[A-Za-z0.9.-]+\.[A-Z|a.z]{2,}\b|Chat Log:?\s*\w+)'

# List of regular expressions and their corresponding categories
checks = [
    (credit_card_regex, 'credit card'),
    (health_record_regex, 'health record'),
    (financial_info_regex, 'financial info'),
    (intellectual_property_regex, 'intellectual property'),
    (legal_documents_regex, 'legal documents'),
    (gov_id_regex, 'governmental ID'),
    (educational_records_regex, 'educational records'),
    (corporate_info_regex, 'corporate info'),
    (employee_data_regex, 'employee data'),
    (customer_data_regex, 'customer data'),
    (communication_records_regex, 'communication records')
]

# Retention periods in days for each category
retention_periods = {
    'credit card': 365,
    'health record': 6 * 365,
    'financial info': 10 * 365,
    'intellectual property': 7 * 365,
    'legal documents': 7 * 365,
    'governmental ID': 365,
    'educational records': 10 * 365,
    'corporate info': 7 * 365,
    'employee data': 7 * 365,
    'customer data': 5 * 365,
    'communication records': 3 * 365
}

# VirusTotal API key
VIRUSTOTAL_API_KEY = '<api key>'

# Helper functions to read various file types
def read_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def read_xlsx(file_path):
    workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    text = []
    for sheet in workbook.sheetnames:
        worksheet = workbook[sheet]
        for row in worksheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text.append(str(cell))
    return "\n".join(text)

def read_pptx(file_path):
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def read_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfFileReader(file)
        for page_num in range(pdf_reader.numPages):
            page = pdf_reader.getPage(page_num)
            text += page.extractText()
    return text

def read_text(file_path):
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
        return file.read()

def calculate_hash(file_path):
    sha256_hash = hashlib.sha256()
    with open(file_path, "rb") as f:
        for byte_block in iter(lambda: f.read(4096), b""):
            sha256_hash.update(byte_block)
    return sha256_hash.hexdigest()

def check_virustotal(hash):
    url = f"https://www.virustotal.com/api/v3/files/{hash}"
    headers = {"Accept": "application/json", "x-apikey": VIRUSTOTAL_API_KEY}
    response = requests.get(url, headers=headers)
    if response.status_code == 200:
        data = response.json()
        malicious_count = data.get('data', {}).get('attributes', {}).get('last_analysis_stats', {}).get('malicious', 0)
        return "Potential Malicious" if malicious_count > 0 else "Clean"
    else:
        return "Hash was not found in virustotal - Maybe Clean"

def perform_regex_checks(content):
    classifications = []
    max_retention_period = 0
    for regex, category in checks:
        if re.search(regex, content):
            classifications.append(category)
            max_retention_period = max(max_retention_period, retention_periods[category])
    return classifications, max_retention_period

def compute_deletion_date(file_path, retention_days):
    creation_time = os.path.getctime(file_path)
    creation_date = datetime.fromtimestamp(creation_time)
    deletion_date = creation_date + timedelta(days=retention_days)
    return deletion_date.strftime('%Y-%m-%d')

def classify_file(file_path):
    try:
        content = ""
        if file_path.endswith('.docx'):
            content = read_docx(file_path)
        elif file_path.endswith('.xlsx'):
            content = read_xlsx(file_path)
        elif file_path.endswith('.pptx'):
            content = read_pptx(file_path)
        elif file_path.endswith('.pdf'):
            content = read_pdf(file_path)
        elif file_path.endswith('.txt'):
            content = read_text(file_path)

        classifications, retention_days = perform_regex_checks(content)
        deletion_date = compute_deletion_date(file_path, retention_days)
        current_date = datetime.now()

        file_hash = calculate_hash(file_path)
        virustotal_result = check_virustotal(file_hash)

        deletion_status = 'Deleted' if current_date > datetime.strptime(deletion_date, '%Y-%m-%d') and virustotal_result == "Malicious" else 'Retain until ' + deletion_date

        return {
            "file_name": os.path.basename(file_path),
            "Classifications": classifications,
            "virus_total_result": virustotal_result,
            "retention_days": retention_days,
            "Deletion Date": deletion_date,
            "Deletion Status": deletion_status,
            "Hash": file_hash
        }
    except Exception as e:
        return {
            "file_name": os.path.basename(file_path),
            "Error": str(e)
        }

def create_and_output_json(file_paths, directory):
    if not os.path.exists(directory):
        print(f"The directory {directory} is not accessible. Switching to backup directory Z:\\")
        directory = "Z:\\"
        
    results = []
    for file_path in file_paths:
        full_path = os.path.join(directory, file_path)
        if os.path.isfile(full_path):
            file_result = classify_file(full_path)
            results.append(file_result)

    # Save the results to a JSON file in the desired directory
    output_path = 'C:\\Users\\vagrant\\Desktop\\' + socket.gethostbyname(socket.gethostname()) + '_Classification.json'
    with open(output_path, 'w') as f:
        json.dump(results, f, indent=4)
    print(f"Results saved to {output_path}")

# Example usage
directory = r'\\VBOXSVR\share'  # Adjust this path as necessary
create_and_output_json(os.listdir(directory), directory)
